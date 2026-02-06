#!/usr/bin/env python3
"""
PPT Generator - 한국어 최적화 미니멀 프레젠테이션 생성

python-pptx를 사용하여 PPTX 파일을 직접 생성합니다.
"""

import argparse
import json
import os
import sys
from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData


# 컬러 팔레트 정의
PALETTES = {
    "sage": {
        "surface": "#f5f5f0",
        "surface_foreground": "#1a1a1a",
        "primary": "#b8c4b8",
        "accent": "#2d2d2d",
        "muted": "#e8e8e3",
        "muted_foreground": "#666666",
        "dark_bg": "#2d2d2d",
        "dark_fg": "#f5f5f0",
    },
    "mono": {
        "surface": "#ffffff",
        "surface_foreground": "#111111",
        "primary": "#f0f0f0",
        "accent": "#111111",
        "muted": "#f5f5f5",
        "muted_foreground": "#666666",
        "dark_bg": "#111111",
        "dark_fg": "#ffffff",
    },
    "navy": {
        "surface": "#f8f9fc",
        "surface_foreground": "#1a1f36",
        "primary": "#dce3f0",
        "accent": "#111111",
        "muted": "#eef1f6",
        "muted_foreground": "#5c6478",
        "dark_bg": "#1a1f36",
        "dark_fg": "#f8f9fc",
    },
}

# 폰트 설정
FONTS = {
    "display": "Noto Serif KR",  # 제목용
    "content": "Pretendard",      # 본문용
    "fallback_display": "맑은 고딕",
    "fallback_content": "맑은 고딕",
}

# 슬라이드 크기 프리셋
SLIDE_SIZES = {
    "16:9": (Inches(13.333), Inches(7.5)),    # 기본 와이드스크린
    "4:3": (Inches(10), Inches(7.5)),          # 전통적 비율
    "1:1": (Inches(10), Inches(10)),           # 정사각형 (모바일/SNS)
    "4:5": (Inches(8), Inches(10)),            # 세로형 (인스타그램)
    "9:16": (Inches(5.625), Inches(10)),       # 세로형 (스토리)
}

# 기본 슬라이드 크기 (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# 여백
MARGIN = Inches(0.75)

# 원라이너 영역 높이
ONELINER_HEIGHT = Inches(0.5)
SEPARATOR_HEIGHT = Inches(0.02)


def hex_to_rgb(hex_color: str) -> RGBColor:
    """HEX 컬러를 RGBColor로 변환"""
    hex_color = hex_color.lstrip('#')
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )


def set_text_frame_properties(text_frame, font_name: str, font_size: int,
                               font_color: str, bold: bool = False,
                               alignment: PP_ALIGN = PP_ALIGN.LEFT):
    """텍스트 프레임 속성 설정"""
    text_frame.word_wrap = True
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = alignment
        paragraph.font.name = font_name
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = hex_to_rgb(font_color)
        paragraph.font.bold = bold


def add_text_box(slide, left: float, top: float, width: float, height: float,
                 text: str, font_name: str, font_size: int, font_color: str,
                 bold: bool = False, alignment: PP_ALIGN = PP_ALIGN.LEFT,
                 vertical_anchor: MSO_ANCHOR = MSO_ANCHOR.TOP) -> None:
    """텍스트 박스 추가"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = text
    tf.paragraphs[0].alignment = alignment
    tf.paragraphs[0].font.name = font_name
    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].font.color.rgb = hex_to_rgb(font_color)
    tf.paragraphs[0].font.bold = bold

    # 수직 정렬
    txBox.text_frame.auto_size = None
    txBox.text_frame.word_wrap = True


def add_shape_with_fill(slide, left: float, top: float, width: float,
                        height: float, fill_color: str) -> None:
    """배경 사각형 추가"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
    shape.line.fill.background()


def get_dimensions():
    """현재 슬라이드 크기에서 인치 단위 값 반환"""
    width = SLIDE_WIDTH.inches
    height = SLIDE_HEIGHT.inches
    margin = MARGIN.inches
    content_width = width - (2 * margin)
    return width, height, margin, content_width


def add_oneliner_with_separator(slide, data: Dict, palette: Dict, margin: float, content_width: float) -> float:
    """원라이너와 구분선을 슬라이드 최상단에 추가하고, 콘텐츠 시작 Y 위치 반환"""
    oneliner = data.get("oneliner", "")

    if not oneliner:
        # 원라이너가 없으면 기존 margin 위치부터 시작
        return margin

    # 원라이너 텍스트 (최상단) - 핵심 메시지이므로 크게, 강조색
    oneliner_y = margin * 0.5
    add_text_box(
        slide, margin, oneliner_y, content_width, 1.0,
        oneliner,
        FONTS["content"], 32, palette["accent"],
        alignment=PP_ALIGN.LEFT
    )

    # 구분선 (원라이너 아래) - 더 두껍게
    separator_y = oneliner_y + 1.1
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(margin), Inches(separator_y),
        Inches(content_width), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = hex_to_rgb(palette["muted"])
    line.line.fill.background()

    # 콘텐츠 시작 Y 위치 반환 (구분선 아래)
    return separator_y + 0.2


def add_footnotes(slide, data: Dict, palette: Dict, margin: float, content_width: float, height: float) -> None:
    """슬라이드 하단 오른쪽에 각주 추가 (*, ** 등의 표기법)"""
    footnotes = data.get("footnotes", [])

    if not footnotes:
        return

    # 각주 텍스트 조합 (예: "* BF16: Brain Float 16bit  ** FP8: 8-bit Floating Point")
    footnote_text = "  ".join(footnotes)

    # 슬라이드 하단 오른쪽에 배치
    footnote_y = height - margin * 0.6
    add_text_box(
        slide, margin, footnote_y, content_width, 0.4,
        footnote_text,
        FONTS["content"], 10, palette["muted_foreground"],
        alignment=PP_ALIGN.RIGHT
    )


def create_cover_slide(prs: Presentation, data: Dict, palette: Dict) -> None:
    """표지 슬라이드 생성"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    width, height, margin, content_width = get_dimensions()

    # 배경색
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = hex_to_rgb(palette["surface"])
    background.line.fill.background()

    # 제목 (수직 중앙보다 약간 위)
    title_y = height * 0.33
    add_text_box(
        slide, margin, title_y, content_width, 1.5,
        data.get("title", "제목"),
        FONTS["display"], 44, palette["surface_foreground"],
        bold=True, alignment=PP_ALIGN.CENTER
    )

    # 부제목 (제목이 2줄일 수 있으므로 충분한 간격 확보)
    if data.get("subtitle"):
        subtitle_y = height * 0.58
        add_text_box(
            slide, margin, subtitle_y, content_width, 0.8,
            data["subtitle"],
            FONTS["content"], 20, palette["muted_foreground"],
            alignment=PP_ALIGN.CENTER
        )

    # 발표자 / 날짜
    footer_text = ""
    if data.get("author"):
        footer_text = data["author"]
    if data.get("date"):
        footer_text += f"  |  {data['date']}" if footer_text else data["date"]

    if footer_text:
        footer_y = height * 0.85
        add_text_box(
            slide, margin, footer_y, content_width, 0.5,
            footer_text,
            FONTS["content"], 14, palette["muted_foreground"],
            alignment=PP_ALIGN.CENTER
        )


def create_section_slide(prs: Presentation, data: Dict, palette: Dict) -> None:
    """섹션 구분 슬라이드 생성 (다크 배경)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    width, height, margin, content_width = get_dimensions()

    # 다크 배경
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = hex_to_rgb(palette["dark_bg"])
    background.line.fill.background()

    # 배지
    if data.get("badge"):
        badge_y = height * 0.38
        add_text_box(
            slide, margin, badge_y, content_width, 0.5,
            data["badge"],
            FONTS["content"], 14, palette["muted_foreground"],
            alignment=PP_ALIGN.CENTER
        )

    # 섹션 제목
    title_y = height * 0.45
    add_text_box(
        slide, margin, title_y, content_width, 1.2,
        data.get("title", "섹션 제목"),
        FONTS["display"], 36, palette["dark_fg"],
        bold=True, alignment=PP_ALIGN.CENTER
    )


def create_stats_grid_slide(prs: Presentation, data: Dict, palette: Dict) -> None:
    """통계 그리드 슬라이드 생성"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    width, height, margin, content_width = get_dimensions()

    # 배경
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = hex_to_rgb(palette["surface"])
    background.line.fill.background()

    # 원라이너 + 구분선 (최상단)
    content_start_y = add_oneliner_with_separator(slide, data, palette, margin, content_width)

    # 슬라이드 제목
    add_text_box(
        slide, margin, content_start_y, content_width, 0.8,
        data.get("title", "핵심 지표"),
        FONTS["display"], 28, palette["muted_foreground"],
        bold=True
    )

    # 통계 항목들
    stats = data.get("stats", [])
    num_stats = len(stats)

    if num_stats > 0:
        # 그리드 계산 - 상대적 위치
        col_width = content_width / min(num_stats, 4)
        start_y = content_start_y + 1.3

        for i, stat in enumerate(stats[:4]):  # 최대 4개
            col_x = margin + (i * col_width)

            # 숫자
            add_text_box(
                slide, col_x, start_y, col_width - 0.3, 0.8,
                stat.get("number", "0"),
                FONTS["display"], 48, palette["accent"],
                bold=True, alignment=PP_ALIGN.CENTER
            )

            # 단위 (unit) - 숫자 바로 아래
            if stat.get("unit"):
                add_text_box(
                    slide, col_x, start_y + 0.8, col_width - 0.3, 0.4,
                    stat["unit"],
                    FONTS["content"], 18, palette["accent"],
                    alignment=PP_ALIGN.CENTER
                )

            # 명칭 (label) - 단위 아래
            add_text_box(
                slide, col_x, start_y + 1.3, col_width - 0.3, 0.5,
                stat.get("label", ""),
                FONTS["content"], 14, palette["muted_foreground"],
                alignment=PP_ALIGN.CENTER
            )

    # 본문 (body) - 숫자 아래 설명 텍스트
    if data.get("body"):
        body_y = content_start_y + 3.8  # 숫자와 레이블 아래
        add_text_box(
            slide, margin, body_y, content_width, 1.5,
            data["body"],
            FONTS["content"], 16, palette["muted_foreground"],
            alignment=PP_ALIGN.LEFT
        )

    # 각주 (footnotes) - 슬라이드 하단
    add_footnotes(slide, data, palette, margin, content_width, height)


def create_two_column_slide(prs: Presentation, data: Dict, palette: Dict) -> None:
    """2단 비교 슬라이드 생성"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    width, height, margin, content_width = get_dimensions()

    # 배경
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = hex_to_rgb(palette["surface"])
    background.line.fill.background()

    # 원라이너 + 구분선 (최상단)
    content_start_y = add_oneliner_with_separator(slide, data, palette, margin, content_width)

    # 슬라이드 제목
    add_text_box(
        slide, margin, content_start_y, content_width, 0.8,
        data.get("title", "비교 분석"),
        FONTS["display"], 28, palette["muted_foreground"],
        bold=True
    )

    # 2컬럼 계산
    gap = 0.5
    col_width = (content_width - gap) / 2
    left_x = margin
    right_x = margin + col_width + gap
    header_y = content_start_y + 0.9
    items_start_y = content_start_y + 1.5

    # 왼쪽 컬럼
    left_data = data.get("left", {})
    add_text_box(
        slide, left_x, header_y, col_width, 0.6,
        left_data.get("heading", "왼쪽"),
        FONTS["display"], 24, palette["surface_foreground"],
        bold=True
    )

    left_items = left_data.get("items", [])
    for j, item in enumerate(left_items[:5]):
        add_text_box(
            slide, left_x, items_start_y + (j * 0.6), col_width, 0.5,
            f"• {item}",
            FONTS["content"], 16, palette["surface_foreground"]
        )

    # 오른쪽 컬럼
    right_data = data.get("right", {})
    add_text_box(
        slide, right_x, header_y, col_width, 0.6,
        right_data.get("heading", "오른쪽"),
        FONTS["display"], 24, palette["surface_foreground"],
        bold=True
    )

    right_items = right_data.get("items", [])
    for j, item in enumerate(right_items[:5]):
        add_text_box(
            slide, right_x, items_start_y + (j * 0.6), col_width, 0.5,
            f"• {item}",
            FONTS["content"], 16, palette["surface_foreground"]
        )

    # 각주 (footnotes) - 슬라이드 하단
    add_footnotes(slide, data, palette, margin, content_width, height)


def create_three_column_slide(prs: Presentation, data: Dict, palette: Dict) -> None:
    """3단 컬럼 슬라이드 생성"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    width, height, margin, content_width = get_dimensions()

    # 배경
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = hex_to_rgb(palette["surface"])
    background.line.fill.background()

    # 원라이너 + 구분선 (최상단)
    content_start_y = add_oneliner_with_separator(slide, data, palette, margin, content_width)

    # 슬라이드 제목
    add_text_box(
        slide, margin, content_start_y, content_width, 0.8,
        data.get("title", "핵심 기능"),
        FONTS["display"], 28, palette["muted_foreground"],
        bold=True
    )

    columns = data.get("columns", [])
    gap = 0.3
    col_width = (content_width - (2 * gap)) / 3
    header_y = content_start_y + 0.9
    desc_y = content_start_y + 1.5

    for i, col in enumerate(columns[:3]):
        col_x = margin + (i * (col_width + gap))

        # 컬럼 제목
        add_text_box(
            slide, col_x, header_y, col_width, 0.6,
            col.get("heading", f"기능 {i+1}"),
            FONTS["display"], 20, palette["surface_foreground"],
            bold=True
        )

        # 컬럼 설명
        add_text_box(
            slide, col_x, desc_y, col_width, 2.5,
            col.get("description", ""),
            FONTS["content"], 14, palette["muted_foreground"]
        )


def create_image_text_slide(prs: Presentation, data: Dict, palette: Dict) -> None:
    """이미지 + 텍스트 슬라이드 생성"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    width, height, margin, content_width = get_dimensions()

    # 배경
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = hex_to_rgb(palette["surface"])
    background.line.fill.background()

    # 원라이너 + 구분선 (최상단)
    content_start_y = add_oneliner_with_separator(slide, data, palette, margin, content_width)

    # 슬라이드 제목
    add_text_box(
        slide, margin, content_start_y, content_width, 0.8,
        data.get("title", "스토리텔링"),
        FONTS["display"], 28, palette["muted_foreground"],
        bold=True
    )

    # 텍스트가 있는지 확인
    text_content = data.get("text", "").strip()
    has_text = bool(text_content)

    # 텍스트가 없으면 이미지를 전체 너비로, 있으면 2컬럼
    if has_text:
        gap = 0.5
        col_width = (content_width - gap) / 2
        image_width = col_width
    else:
        # 텍스트 없으면 이미지 전체 너비 사용
        image_width = content_width

    content_y = content_start_y + 0.9
    content_height = height - content_y - margin - 0.3  # 각주 공간 확보

    # 이미지 영역
    image_path = data.get("image_path")
    if image_path and os.path.exists(image_path):
        try:
            from PIL import Image as PILImage
            # 이미지 원본 크기 확인
            with PILImage.open(image_path) as img:
                img_width, img_height = img.size

            # 이미지 비율 계산
            img_ratio = img_width / img_height

            # 사용 가능한 영역
            available_width = image_width
            available_height = content_height
            area_ratio = available_width / available_height

            # 비율에 맞게 크기 조정 (영역 안에 맞추기)
            if img_ratio > area_ratio:
                # 이미지가 더 넓음 → 너비 기준
                final_width = available_width
                final_height = available_width / img_ratio
            else:
                # 이미지가 더 높음 → 높이 기준
                final_height = available_height
                final_width = available_height * img_ratio

            slide.shapes.add_picture(
                image_path,
                Inches(margin), Inches(content_y),
                width=Inches(final_width), height=Inches(final_height)
            )
        except Exception as e:
            # 이미지 로드 실패 시 placeholder
            placeholder = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(margin), Inches(content_y), Inches(image_width), Inches(content_height)
            )
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = hex_to_rgb(palette["muted"])
            placeholder.line.fill.background()
    else:
        # 이미지 없을 때 placeholder
        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(margin), Inches(content_y), Inches(image_width), Inches(content_height)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = hex_to_rgb(palette["muted"])
        placeholder.line.fill.background()

    # 텍스트 영역 (텍스트가 있을 때만)
    if has_text:
        gap = 0.5
        col_width = (content_width - gap) / 2
        right_x = margin + col_width + gap
        add_text_box(
            slide, right_x, content_y, col_width, content_height,
            text_content,
            FONTS["content"], 16, palette["surface_foreground"]
        )

    # 각주 (footnotes) - 슬라이드 하단 오른쪽
    add_footnotes(slide, data, palette, margin, content_width, height)


def create_timeline_slide(prs: Presentation, data: Dict, palette: Dict) -> None:
    """타임라인 슬라이드 생성 (단계별 로드맵)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    width, height, margin, content_width = get_dimensions()

    # 배경
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = hex_to_rgb(palette["surface"])
    background.line.fill.background()

    # 원라이너 + 구분선 (최상단)
    content_start_y = add_oneliner_with_separator(slide, data, palette, margin, content_width)

    # 슬라이드 제목
    add_text_box(
        slide, margin, content_start_y, content_width, 0.8,
        data.get("title", "로드맵"),
        FONTS["display"], 28, palette["muted_foreground"],
        bold=True
    )

    # 타임라인 단계들
    steps = data.get("steps", [])
    num_steps = len(steps)

    if num_steps == 0:
        return

    # 타임라인 영역 계산 - 슬라이드 중앙보다 약간 아래에 배치
    timeline_y = height * 0.55  # 상단 텍스트 공간 확보
    step_width = content_width / num_steps
    arrow_height = 0.7

    # 그라데이션 색상 정의 (진한 청록 → 연한 청록)
    gradient_colors = [
        "#0d9488",  # teal-600
        "#14b8a6",  # teal-500
        "#2dd4bf",  # teal-400
        "#5eead4",  # teal-300
        "#99f6e4",  # teal-200
    ]

    for i, step in enumerate(steps[:5]):  # 최대 5단계
        step_x = margin + (i * step_width)
        center_x = step_x + (step_width / 2)

        # 화살표 도형 (CHEVRON)
        arrow_left = step_x + 0.1
        arrow_width = step_width - 0.2

        chevron = slide.shapes.add_shape(
            MSO_SHAPE.CHEVRON,
            Inches(arrow_left), Inches(timeline_y),
            Inches(arrow_width), Inches(arrow_height)
        )
        chevron.fill.solid()
        color_idx = min(i, len(gradient_colors) - 1)
        chevron.fill.fore_color.rgb = hex_to_rgb(gradient_colors[color_idx])
        chevron.line.fill.background()

        # 단계 번호 (화살표 안)
        step_num = step.get("number", f"{i+1}")
        add_text_box(
            slide, arrow_left + 0.2, timeline_y + 0.15, arrow_width - 0.4, 0.4,
            step_num,
            FONTS["content"], 14, "#ffffff",
            bold=True, alignment=PP_ALIGN.CENTER
        )

        # 원형 아이콘 (화살표 위 또는 아래)
        is_top = (i % 2 == 0)  # 홀짝으로 위/아래 교대
        circle_size = 0.5
        circle_x = center_x - (circle_size / 2)

        if is_top:
            circle_y = timeline_y - circle_size - 0.15
        else:
            circle_y = timeline_y + arrow_height + 0.15

        # 원형 도형
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(circle_x), Inches(circle_y),
            Inches(circle_size), Inches(circle_size)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = hex_to_rgb(gradient_colors[color_idx])
        circle.line.color.rgb = hex_to_rgb("#ffffff")
        circle.line.width = Pt(2)

        # 연결선 (원과 화살표 연결)
        if is_top:
            line_start_y = circle_y + circle_size
            line_end_y = timeline_y
        else:
            line_start_y = timeline_y + arrow_height
            line_end_y = circle_y

        connector = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(center_x - 0.015), Inches(line_start_y),
            Inches(0.03), Inches(abs(line_end_y - line_start_y))
        )
        connector.fill.solid()
        connector.fill.fore_color.rgb = hex_to_rgb(gradient_colors[color_idx])
        connector.line.fill.background()

        # 단계 제목과 설명 위치 계산
        title_text = step.get("title", f"단계 {i+1}")
        desc_text = step.get("description", "")

        if is_top:
            # 위쪽 배치: 원 위에 텍스트 (원에 가깝게)
            title_y = circle_y - 0.45
            desc_y = circle_y - 0.75
        else:
            # 아래쪽 배치: 원 아래에 텍스트
            title_y = circle_y + circle_size + 0.1
            desc_y = title_y + 0.4

        # 단계 제목
        add_text_box(
            slide, step_x, title_y, step_width, 0.4,
            title_text,
            FONTS["display"], 14, palette["surface_foreground"],
            bold=True, alignment=PP_ALIGN.CENTER
        )

        # 단계 설명
        if desc_text:
            add_text_box(
                slide, step_x, desc_y, step_width, 0.5,
                desc_text,
                FONTS["content"], 10, palette["muted_foreground"],
                alignment=PP_ALIGN.CENTER
            )

    # 각주 (footnotes) - 슬라이드 하단
    add_footnotes(slide, data, palette, margin, content_width, height)


def create_bar_chart_slide(prs: Presentation, data: Dict, palette: Dict) -> None:
    """막대 그래프 슬라이드 생성 (도형으로 직접 그리기)

    가이드 원칙:
    - 자동차트 대신 도형으로 직접 그리기
    - Y축 기준선 생략, 심플하게
    - 강조 항목은 색상으로 차별화
    - 필수: 단위, 출처, 항목 이름
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    width, height, margin, content_width = get_dimensions()

    # 배경
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = hex_to_rgb(palette["surface"])
    background.line.fill.background()

    # 원라이너 + 구분선 (최상단)
    content_start_y = add_oneliner_with_separator(slide, data, palette, margin, content_width)

    # 슬라이드 제목
    add_text_box(
        slide, margin, content_start_y, content_width, 0.8,
        data.get("title", "비교"),
        FONTS["display"], 28, palette["muted_foreground"],
        bold=True
    )

    # 차트 데이터
    bars = data.get("bars", [])
    unit = data.get("unit", "")
    source = data.get("source", "")
    highlight_index = data.get("highlight", -1)  # 강조할 막대 인덱스 (0부터)

    if not bars:
        return

    num_bars = len(bars)

    # 차트 영역 계산
    chart_area_y = content_start_y + 1.2
    chart_area_height = 4.0
    chart_area_x = margin + 1.5
    chart_area_width = content_width - 3.0

    # 막대 설정
    bar_gap = 0.8  # 막대 간 간격
    total_gap = bar_gap * (num_bars - 1)
    bar_width = (chart_area_width - total_gap) / num_bars
    bar_width = min(bar_width, 2.5)  # 최대 너비 제한

    # 전체 막대 영역 중앙 정렬
    total_bars_width = (bar_width * num_bars) + total_gap
    start_x = chart_area_x + (chart_area_width - total_bars_width) / 2

    # 최대값 찾기 (스케일 계산용)
    max_value = max(bar.get("value", 0) for bar in bars)

    # 색상 정의
    default_color = "#9ca3af"  # gray-400
    highlight_color = "#14b8a6"  # teal-500

    for i, bar in enumerate(bars[:5]):  # 최대 5개
        value = bar.get("value", 0)
        label = bar.get("label", "")

        # 막대 높이 계산 (비례)
        bar_height = (value / max_value) * (chart_area_height - 1.0) if max_value > 0 else 0
        bar_x = start_x + i * (bar_width + bar_gap)
        bar_y = chart_area_y + (chart_area_height - bar_height)

        # 막대 색상 (강조 여부)
        is_highlight = (i == highlight_index)
        bar_color = highlight_color if is_highlight else default_color

        # 막대 도형
        bar_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(bar_x), Inches(bar_y),
            Inches(bar_width), Inches(bar_height)
        )
        bar_shape.fill.solid()
        bar_shape.fill.fore_color.rgb = hex_to_rgb(bar_color)
        bar_shape.line.fill.background()

        # 데이터 라벨 (막대 위)
        label_color = highlight_color if is_highlight else palette["surface_foreground"]
        value_text = f"{value:,}" if isinstance(value, int) else str(value)
        add_text_box(
            slide, bar_x, bar_y - 0.5, bar_width, 0.4,
            value_text,
            FONTS["content"], 16 if is_highlight else 14,
            label_color,
            bold=is_highlight, alignment=PP_ALIGN.CENTER
        )

        # 항목 라벨 (막대 아래)
        label_y = chart_area_y + chart_area_height + 0.1
        add_text_box(
            slide, bar_x, label_y, bar_width, 0.4,
            label,
            FONTS["content"], 12 if is_highlight else 11,
            palette["surface_foreground"] if is_highlight else palette["muted_foreground"],
            bold=is_highlight, alignment=PP_ALIGN.CENTER
        )

    # 단위 표시 (우측 상단)
    if unit:
        add_text_box(
            slide, margin + content_width - 2.0, chart_area_y - 0.3, 2.0, 0.3,
            f"(단위: {unit})",
            FONTS["content"], 10, palette["muted_foreground"],
            alignment=PP_ALIGN.RIGHT
        )

    # 출처 표시 (우측 하단)
    if source:
        add_text_box(
            slide, margin + content_width - 3.0, chart_area_y + chart_area_height + 0.6, 3.0, 0.3,
            f"출처: {source}",
            FONTS["content"], 9, palette["muted_foreground"],
            alignment=PP_ALIGN.RIGHT
        )

    # 각주 (footnotes)
    add_footnotes(slide, data, palette, margin, content_width, height)


def create_pie_chart_slide(prs: Presentation, data: Dict, palette: Dict) -> None:
    """원형 그래프(도넛 차트) 슬라이드 생성

    가이드 원칙:
    - 파이 조각 5개 이내 (초과 시 "기타"로 묶기)
    - 12시 정각 기준 오른쪽부터 시계방향 배치
    - 중요 조각은 색상으로 강조
    - 필수: 단위, 출처, 제목, 항목 이름
    - 강조 항목은 중앙에 라벨+값 크게 표시
    """
    import math

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    width, height, margin, content_width = get_dimensions()

    # 배경
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = hex_to_rgb(palette["surface"])
    background.line.fill.background()

    # 원라이너 + 구분선 (최상단)
    content_start_y = add_oneliner_with_separator(slide, data, palette, margin, content_width)

    # 슬라이드 제목
    add_text_box(
        slide, margin, content_start_y, content_width, 0.8,
        data.get("title", "비율"),
        FONTS["display"], 28, palette["muted_foreground"],
        bold=True
    )

    # 차트 데이터
    slices = data.get("slices", [])
    unit = data.get("unit", "%")
    source = data.get("source", "")
    highlight_index = data.get("highlight", 0)  # 강조할 조각 인덱스 (기본: 첫 번째)

    if not slices:
        return

    # 최대 5개로 제한
    slices = slices[:5]
    num_slices = len(slices)

    # 전체 합계 계산
    total = sum(slice_data.get("value", 0) for slice_data in slices)
    if total == 0:
        return

    # 차트 영역 계산
    chart_center_x = margin + content_width * 0.4  # 왼쪽에 배치
    chart_center_y = content_start_y + 2.8
    outer_radius = 1.8
    inner_radius = 1.0  # 도넛 내부 반지름

    # 색상 팔레트 (진한 청록 → 연한 회색 그라데이션)
    pie_colors = [
        "#0f766e",  # teal-700 (강조용)
        "#14b8a6",  # teal-500
        "#5eead4",  # teal-300
        "#99f6e4",  # teal-200
        "#ccfbf1",  # teal-100
    ]
    muted_colors = [
        "#374151",  # gray-700
        "#6b7280",  # gray-500
        "#9ca3af",  # gray-400
        "#d1d5db",  # gray-300
        "#e5e7eb",  # gray-200
    ]

    # 파이 조각 그리기 (MSO_SHAPE.BLOCK_ARC 대신 python-pptx 차트 기능 사용)
    chart_data = CategoryChartData()
    chart_data.categories = [s.get("label", "") for s in slices]
    chart_data.add_series('', [s.get("value", 0) for s in slices])

    chart_x = chart_center_x - outer_radius
    chart_y = chart_center_y - outer_radius
    chart_size = outer_radius * 2

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        Inches(chart_x), Inches(chart_y),
        Inches(chart_size), Inches(chart_size),
        chart_data
    )
    chart = chart_shape.chart

    # 차트 스타일링 (미니멀하게)
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = False

    # 각 조각 색상 설정
    series = chart.series[0]
    for i, point in enumerate(series.points):
        point.format.fill.solid()
        if i == highlight_index:
            point.format.fill.fore_color.rgb = hex_to_rgb(pie_colors[0])
        else:
            # 비강조 조각은 연한 색상
            color_idx = min(i + 1, len(muted_colors) - 1)
            point.format.fill.fore_color.rgb = hex_to_rgb(muted_colors[color_idx])

    # 강조 항목 라벨 (중앙)
    if 0 <= highlight_index < num_slices:
        highlight_slice = slices[highlight_index]
        highlight_label = highlight_slice.get("label", "")
        highlight_value = highlight_slice.get("value", 0)

        # 라벨 (작게)
        add_text_box(
            slide, chart_center_x - 0.8, chart_center_y - 0.5, 1.6, 0.4,
            highlight_label,
            FONTS["content"], 14, pie_colors[0],
            bold=False, alignment=PP_ALIGN.CENTER
        )
        # 값 (크게)
        value_text = f"{highlight_value}{unit}"
        add_text_box(
            slide, chart_center_x - 1.0, chart_center_y - 0.1, 2.0, 0.6,
            value_text,
            FONTS["display"], 32, pie_colors[0],
            bold=True, alignment=PP_ALIGN.CENTER
        )

    # 오른쪽에 범례 (수동 생성)
    legend_x = margin + content_width * 0.65
    legend_y = content_start_y + 1.5
    legend_item_height = 0.5

    for i, slice_data in enumerate(slices):
        label = slice_data.get("label", "")
        value = slice_data.get("value", 0)
        is_highlight = (i == highlight_index)

        # 색상 박스
        color_box_size = 0.2
        if is_highlight:
            box_color = pie_colors[0]
        else:
            color_idx = min(i + 1, len(muted_colors) - 1)
            box_color = muted_colors[color_idx]

        color_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(legend_x), Inches(legend_y + i * legend_item_height + 0.1),
            Inches(color_box_size), Inches(color_box_size)
        )
        color_box.fill.solid()
        color_box.fill.fore_color.rgb = hex_to_rgb(box_color)
        color_box.line.fill.background()

        # 라벨
        label_text = f"{label}"
        add_text_box(
            slide, legend_x + 0.35, legend_y + i * legend_item_height, 1.5, 0.4,
            label_text,
            FONTS["content"], 12 if is_highlight else 11,
            palette["surface_foreground"] if is_highlight else palette["muted_foreground"],
            bold=is_highlight
        )

        # 값
        value_text = f"{value:,}{unit}" if isinstance(value, int) else f"{value}{unit}"
        add_text_box(
            slide, legend_x + 1.9, legend_y + i * legend_item_height, 1.0, 0.4,
            value_text,
            FONTS["content"], 12 if is_highlight else 11,
            palette["surface_foreground"] if is_highlight else palette["muted_foreground"],
            bold=is_highlight, alignment=PP_ALIGN.RIGHT
        )

    # 단위 표시 (우측 상단)
    add_text_box(
        slide, margin + content_width - 2.0, content_start_y + 0.9, 2.0, 0.3,
        f"(단위: {unit})",
        FONTS["content"], 10, palette["muted_foreground"],
        alignment=PP_ALIGN.RIGHT
    )

    # 출처 표시 (좌측 하단)
    if source:
        add_text_box(
            slide, margin, height - margin - 0.3, 3.0, 0.3,
            f"출처: {source}",
            FONTS["content"], 9, palette["muted_foreground"]
        )

    # 각주 (footnotes)
    add_footnotes(slide, data, palette, margin, content_width, height)


def create_table_slide(prs: Presentation, data: Dict, palette: Dict) -> None:
    """테이블(표) 슬라이드 생성

    가이드 원칙:
    1. 모든 선 0.5pt (색상: rgb 127,127,127)
    2. 양 옆 선 제거
    3. 상단/왼쪽 항목 굵게
    4. 상단 테두리선 1.5pt
    5. 헤더 배경 연한 색
    6. 강조 행 하이라이트
    7. 필수 정보 (제목, 단위, 출처)
    """
    from pptx.util import Pt as PtUtil
    from pptx.oxml.ns import nsdecls
    from pptx.oxml import parse_xml

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    width, height, margin, content_width = get_dimensions()

    # 배경
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = hex_to_rgb(palette["surface"])
    background.line.fill.background()

    # 원라이너 + 구분선 (최상단)
    content_start_y = add_oneliner_with_separator(slide, data, palette, margin, content_width)

    # 슬라이드 제목
    add_text_box(
        slide, margin, content_start_y, content_width - 2.5, 0.8,
        data.get("title", "비교 테이블"),
        FONTS["display"], 28, palette["muted_foreground"],
        bold=True
    )

    # 단위 표시 (우측 상단)
    unit = data.get("unit", "")
    if unit:
        add_text_box(
            slide, margin + content_width - 2.0, content_start_y + 0.2, 2.0, 0.3,
            f"(단위: {unit})",
            FONTS["content"], 10, palette["muted_foreground"],
            alignment=PP_ALIGN.RIGHT
        )

    # 테이블 데이터
    headers = data.get("headers", [])  # ["", "항목1", "항목2", ...]
    rows = data.get("rows", [])  # [{"label": "행1", "values": [1, 2, ...]}, ...]
    highlight_row = data.get("highlight_row", -1)  # 강조할 행 인덱스

    if not headers or not rows:
        return

    num_cols = len(headers)
    num_rows = len(rows) + 1  # 헤더 포함

    # 테이블 영역 계산
    table_y = content_start_y + 1.0
    table_height = min(len(rows) * 0.6 + 0.7, 4.5)  # 최대 높이 제한
    table_width = content_width - 1.0

    # 테이블 생성
    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(margin + 0.5), Inches(table_y),
        Inches(table_width), Inches(table_height)
    )
    table = table_shape.table

    # 열 너비 설정 (첫 번째 열은 좁게)
    first_col_width = table_width * 0.18
    other_col_width = (table_width - first_col_width) / (num_cols - 1) if num_cols > 1 else table_width
    table.columns[0].width = Inches(first_col_width)
    for i in range(1, num_cols):
        table.columns[i].width = Inches(other_col_width)

    # 색상 정의
    line_color = RGBColor(127, 127, 127)  # 선 색상
    header_bg = hex_to_rgb("#e8eaed")  # 헤더 배경 (연한 회색)
    highlight_bg = hex_to_rgb("#d1fae5")  # 강조 배경 (연한 민트)
    text_color = hex_to_rgb(palette["surface_foreground"])
    muted_color = hex_to_rgb(palette["muted_foreground"])

    def set_cell_border(cell, top=None, bottom=None, left=None, right=None):
        """셀 테두리 설정"""
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()

        # 기존 테두리 제거
        for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
            existing = tcPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}' + border_name)
            if existing is not None:
                tcPr.remove(existing)

        def add_border(name, width_pt, color):
            if width_pt is None:
                # 테두리 제거 (투명)
                border = parse_xml(
                    f'<a:{name} xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" w="0" cap="flat" cmpd="sng" algn="ctr">'
                    f'<a:noFill/></a:{name}>'
                )
            else:
                # 테두리 설정
                width_emu = int(width_pt * 12700)  # pt to EMU
                r, g, b = color.r if hasattr(color, 'r') else 127, color.g if hasattr(color, 'g') else 127, color.b if hasattr(color, 'b') else 127
                border = parse_xml(
                    f'<a:{name} xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" w="{width_emu}" cap="flat" cmpd="sng" algn="ctr">'
                    f'<a:solidFill><a:srgbClr val="{r:02X}{g:02X}{b:02X}"/></a:solidFill>'
                    f'<a:prstDash val="solid"/></a:{name}>'
                )
            tcPr.append(border)

        # 각 테두리 설정
        add_border('lnL', left, line_color)
        add_border('lnR', right, line_color)
        add_border('lnT', top, line_color)
        add_border('lnB', bottom, line_color)

    # 헤더 행 설정
    for col_idx, header_text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = str(header_text)

        # 헤더 배경색
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg

        # 텍스트 스타일
        para = cell.text_frame.paragraphs[0]
        para.font.name = FONTS["content"]
        para.font.size = Pt(11)
        para.font.bold = True
        para.font.color.rgb = text_color
        para.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT

        # 수직 정렬
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 테두리 설정: 상단 1.5pt, 하단 0.5pt, 좌우 없음
        set_cell_border(cell, top=1.5, bottom=0.5, left=None, right=None)

    # 데이터 행 설정
    for row_idx, row_data in enumerate(rows):
        actual_row = row_idx + 1
        is_highlight = (row_idx == highlight_row)

        # 첫 번째 열 (라벨)
        label_cell = table.cell(actual_row, 0)
        label_cell.text = str(row_data.get("label", ""))

        if is_highlight:
            label_cell.fill.solid()
            label_cell.fill.fore_color.rgb = highlight_bg
        else:
            label_cell.fill.background()

        label_para = label_cell.text_frame.paragraphs[0]
        label_para.font.name = FONTS["content"]
        label_para.font.size = Pt(11)
        label_para.font.bold = True
        label_para.font.color.rgb = text_color
        label_para.alignment = PP_ALIGN.LEFT
        label_cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 테두리: 하단만 0.5pt
        set_cell_border(label_cell, top=None, bottom=0.5, left=None, right=None)

        # 값 열들
        values = row_data.get("values", [])
        for col_idx, value in enumerate(values):
            if col_idx + 1 >= num_cols:
                break
            cell = table.cell(actual_row, col_idx + 1)
            cell.text = str(value) if value is not None else ""

            if is_highlight:
                cell.fill.solid()
                cell.fill.fore_color.rgb = highlight_bg
            else:
                cell.fill.background()

            para = cell.text_frame.paragraphs[0]
            para.font.name = FONTS["content"]
            para.font.size = Pt(11)
            para.font.bold = False
            para.font.color.rgb = text_color if is_highlight else muted_color
            para.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # 테두리: 하단만 0.5pt
            set_cell_border(cell, top=None, bottom=0.5, left=None, right=None)

    # 출처 표시 (우측 하단)
    source = data.get("source", "")
    if source:
        source_y = table_y + table_height + 0.2
        add_text_box(
            slide, margin + content_width - 3.0, source_y, 3.0, 0.3,
            f"출처: {source}",
            FONTS["content"], 9, palette["muted_foreground"],
            alignment=PP_ALIGN.RIGHT
        )

    # 각주 (footnotes)
    add_footnotes(slide, data, palette, margin, content_width, height)


def create_closing_slide(prs: Presentation, data: Dict, palette: Dict) -> None:
    """마무리 슬라이드 생성"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    width, height, margin, content_width = get_dimensions()

    # 배경
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = hex_to_rgb(palette["surface"])
    background.line.fill.background()

    # 감사 메시지 (수직 중앙)
    title_y = height * 0.38
    add_text_box(
        slide, margin, title_y, content_width, 1.2,
        data.get("title", "감사합니다"),
        FONTS["display"], 44, palette["surface_foreground"],
        bold=True, alignment=PP_ALIGN.CENTER
    )

    # 연락처
    if data.get("contact"):
        contact_y = height * 0.58
        add_text_box(
            slide, margin, contact_y, content_width, 0.6,
            data["contact"],
            FONTS["content"], 16, palette["muted_foreground"],
            alignment=PP_ALIGN.CENTER
        )


def create_content_slide(prs: Presentation, data: Dict, palette: Dict) -> None:
    """일반 콘텐츠 슬라이드 생성 (기본 레이아웃)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    width, height, margin, content_width = get_dimensions()

    # 배경
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = hex_to_rgb(palette["surface"])
    background.line.fill.background()

    # 원라이너 + 구분선 (최상단)
    content_start_y = add_oneliner_with_separator(slide, data, palette, margin, content_width)

    # 슬라이드 제목
    add_text_box(
        slide, margin, content_start_y, content_width, 0.8,
        data.get("title", "제목"),
        FONTS["display"], 28, palette["muted_foreground"],
        bold=True
    )

    # 본문 내용
    content = data.get("content", "")
    if isinstance(content, list):
        content = "\n".join([f"• {item}" for item in content])

    content_y = content_start_y + 0.9
    content_height = height - content_y - margin
    add_text_box(
        slide, margin, content_y, content_width, content_height,
        content,
        FONTS["content"], 16, palette["surface_foreground"]
    )


# 레이아웃 타입별 생성 함수 매핑
LAYOUT_CREATORS = {
    "cover": create_cover_slide,
    "section": create_section_slide,
    "stats_grid": create_stats_grid_slide,
    "two_column": create_two_column_slide,
    "three_column": create_three_column_slide,
    "image_text": create_image_text_slide,
    "timeline": create_timeline_slide,
    "bar_chart": create_bar_chart_slide,
    "pie_chart": create_pie_chart_slide,
    "table": create_table_slide,
    "closing": create_closing_slide,
    "content": create_content_slide,
}


def generate_pptx(config: Dict, output_path: str, palette_name: str = "sage",
                  size: str = "16:9") -> str:
    """PPTX 파일 생성 메인 함수"""
    global SLIDE_WIDTH, SLIDE_HEIGHT

    # 팔레트 선택
    palette = PALETTES.get(palette_name, PALETTES["sage"])

    # 슬라이드 크기 설정
    if size in SLIDE_SIZES:
        SLIDE_WIDTH, SLIDE_HEIGHT = SLIDE_SIZES[size]
    else:
        SLIDE_WIDTH, SLIDE_HEIGHT = SLIDE_SIZES["16:9"]

    # 프레젠테이션 생성
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 슬라이드 생성
    slides_data = config.get("slides", [])

    for slide_data in slides_data:
        layout = slide_data.get("layout", "content")
        creator_func = LAYOUT_CREATORS.get(layout, create_content_slide)
        creator_func(prs, slide_data, palette)

    # 파일 저장
    output_path = os.path.expanduser(output_path)
    prs.save(output_path)

    return output_path


def main():
    parser = argparse.ArgumentParser(
        description="한국어 최적화 미니멀 PPTX 생성기"
    )
    parser.add_argument(
        "--config", "-c",
        required=True,
        help="슬라이드 설정 JSON 파일 경로"
    )
    parser.add_argument(
        "--output", "-o",
        default="output.pptx",
        help="출력 PPTX 파일 경로 (기본: output.pptx)"
    )
    parser.add_argument(
        "--palette", "-p",
        choices=["sage", "mono", "navy"],
        default="sage",
        help="컬러 팔레트 (기본: sage)"
    )
    parser.add_argument(
        "--size", "-s",
        choices=["16:9", "4:3", "1:1", "4:5", "9:16"],
        default="16:9",
        help="슬라이드 비율 (기본: 16:9, 모바일: 1:1 또는 4:5)"
    )

    args = parser.parse_args()

    # JSON 설정 파일 로드
    config_path = os.path.expanduser(args.config)
    if not os.path.exists(config_path):
        print(f"Error: Config file not found: {config_path}", file=sys.stderr)
        sys.exit(1)

    with open(config_path, 'r', encoding='utf-8') as f:
        config = json.load(f)

    # 팔레트 오버라이드 (config에서 지정된 경우)
    palette_name = config.get("palette", args.palette)

    # 사이즈 오버라이드 (config에서 지정된 경우)
    size = config.get("size", args.size)

    # PPTX 생성
    try:
        output_path = generate_pptx(config, args.output, palette_name, size)
        print(f"PPTX 생성 완료: {output_path}")
        print(f"슬라이드 수: {len(config.get('slides', []))}")
        print(f"팔레트: {palette_name}")
        print(f"사이즈: {size}")
    except Exception as e:
        print(f"Error: PPTX 생성 실패: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
